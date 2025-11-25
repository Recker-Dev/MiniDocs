from docx import Document
from docx.shared import Pt as DPt
from pptx import Presentation
from pptx.util import Pt as PPt

def build_docx(sections, buffer):
    doc = Document()

    for section in sections:
        # Heading with bigger font
        heading = doc.add_heading(section.section_name, level=1)
        for run in heading.runs:
            run.font.size = DPt(24)

        # Spacing below heading
        doc.add_paragraph("")  

        # Body with bigger font
        p = doc.add_paragraph(section.content)
        for run in p.runs:
            run.font.size = DPt(14)

        doc.add_page_break()

    doc.save(buffer)
    buffer.seek(0)


def build_pptx(sections, buffer):
    prs = Presentation()

    for section in sections:
        slide_layout = prs.slide_layouts[1]  # Title + content
        slide = prs.slides.add_slide(slide_layout)

        # Title
        title = slide.shapes.title
        title.text = section.section_name
        title.text_frame.paragraphs[0].font.size = PPt(24)

        # Body
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()  # Remove default
        p = tf.add_paragraph()
        p.text = section.content
        p.font.size = PPt(14)

        # Add spacing from title
        tf.margin_top = PPt(12)

    prs.save(buffer)
    buffer.seek(0)
